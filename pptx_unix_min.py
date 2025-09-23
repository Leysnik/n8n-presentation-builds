#!/usr/bin/env python3
# -*- coding: utf-8 -*-

import base64
import subprocess
import shutil
from pathlib import Path

import requests
from pptx import Presentation

# Харды
PPTX_PATH = Path("presentation_demo.pptx").resolve()
API_URL = "http://localhost:5678/webhook-test/webhook-test/pptx"
OUT_DIR = Path("./_pptx_out").resolve()
DPI = 144  # можно поднять до 200–300 для более четких PNG


def run(cmd, cwd=None):
    proc = subprocess.run(cmd, cwd=cwd, stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True)
    if proc.returncode != 0:
        raise RuntimeError(f"Command failed: {' '.join(cmd)}\n{proc.stderr}")
    return proc


def extract_notes(pptx_path: Path):
    prs = Presentation(str(pptx_path))
    notes = []
    for slide in prs.slides:
        txt = ""
        try:
            tf = slide.notes_slide.notes_text_frame
            txt = (tf.text or "").strip() if tf else ""
        except Exception:
            txt = ""
        notes.append(txt)
    return notes


def export_images_unix(pptx_path: Path, out_dir: Path, dpi: int = 144):
    out_dir.mkdir(parents=True, exist_ok=True)

    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    if not soffice:
        raise RuntimeError("LibreOffice (soffice) not found in PATH")

    # 1) PPTX -> PDF
    run([soffice, "--headless", "--convert-to", "pdf", "--outdir", str(out_dir), str(pptx_path)])
    pdf = next(out_dir.glob("*.pdf"), None)
    if not pdf:
        raise RuntimeError("PDF not produced by LibreOffice")

    # 2) PDF -> PNGs
    gs = shutil.which("gs")
    if not gs:
        raise RuntimeError("Ghostscript (gs) not found in PATH")

    pattern = str(out_dir / "slide-%03d.png")
    run([gs, "-sDEVICE=pngalpha", "-o", pattern, f"-r{dpi}", str(pdf)])

    images = sorted(out_dir.glob("slide-*.png"))
    if not images:
        raise RuntimeError("No PNG slides produced by Ghostscript")
    return images


def b64_file(p: Path) -> str:
    return base64.b64encode(p.read_bytes()).decode("ascii")


def main():
    if not PPTX_PATH.exists():
        raise FileNotFoundError(f"{PPTX_PATH} not found")

    notes = extract_notes(PPTX_PATH)
    images = export_images_unix(PPTX_PATH, OUT_DIR, dpi=DPI)

    # Выравниваем по порядку; если расхождение в количестве — берем минимум
    n = min(len(notes), len(images))
    slides = []
    for i in range(n):
        slides.append({
            "slide_index": i + 1,
            "notes": notes[i],
            "image_base64": b64_file(images[i]),
            "filename": PPTX_PATH.name,
            "total_slides": len(notes),
        })

    # Отправляем одним батчем
    payload = {"slides": slides}
    resp = requests.post(API_URL, json=payload, timeout=3600)
    resp.raise_for_status()
    with open('video.mp4', 'wb') as f:
        f.write(resp.content)
    print(f"OK: sent {len(slides)} slide(s) to {API_URL}")


if __name__ == "__main__":
    main()
