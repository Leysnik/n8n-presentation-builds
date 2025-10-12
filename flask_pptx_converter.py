from flask import Flask, request, jsonify
from pptx import Presentation
import base64
import subprocess
import shutil
from pathlib import Path
import tempfile
import uuid

app = Flask(__name__)

@app.route('/convert', methods=['POST'])
def convert_presentation():
    print('//////////////')
    print(request.files)
    print(request.__dict__)
    if 'file' not in request.files:
        return jsonify({'error': 'Файл не найден в запросе'}), 400
    file = request.files['file']

    if file.filename == '' or not file.filename.lower().endswith('.pptx'):
        return jsonify({'error': 'Ожидался файл .pptx'}), 400

    with tempfile.TemporaryDirectory() as temp_dir:
        temp_path = Path(temp_dir)

        unique_id = uuid.uuid4().hex
        pptx_filename = f"upload_{unique_id}.pptx"
        pptx_path = temp_path / pptx_filename
        file.save(str(pptx_path))

        try:
            prs = Presentation(str(pptx_path))
        except Exception as e:
            return jsonify({'error': f'Не удалось открыть презентацию: {str(e)}'}), 400

        notes = []
        for slide in prs.slides:
            txt = ''
            try:
                if slide.has_notes_slide:
                    tf = slide.notes_slide.notes_text_frame
                    txt = (tf.text or '').strip() if tf else ''
            except Exception:
                pass
            notes.append(txt)

        pdf_output = temp_path / f"output_{unique_id}.pdf"

        soffice = shutil.which('soffice') or shutil.which('libreoffice')
        if not soffice:
            return jsonify({'error': 'LibreOffice не найден. Установите и добавьте в PATH.'}), 500

        try:
            subprocess.run([
                soffice, '--headless', '--convert-to', 'pdf',
                '--outdir', str(temp_path), str(pptx_path)
            ], check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
        except subprocess.CalledProcessError:
            return jsonify({'error': 'Ошибка при конвертации в PDF'}), 500

        pdf_path = next(temp_path.glob('*.pdf'), None)
        if not pdf_path:
            return jsonify({'error': 'PDF не был создан'}), 500

        png_pattern = str(temp_path / f"slide_{unique_id}-%03d.png")
        gs = shutil.which('gs')
        if not gs:
            return jsonify({'error': 'Ghostscript (gs) не найден. Установите и добавьте в PATH.'}), 500

        try:
            subprocess.run([
                gs, '-sDEVICE=pngalpha', '-o', png_pattern, '-r144', str(pdf_path)
            ], check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
        except subprocess.CalledProcessError:
            return jsonify({'error': 'Ошибка при конвертации PDF в PNG'}), 500

        image_paths = sorted(temp_path.glob(f"slide_{unique_id}-*.png"))
        if not image_paths:
            return jsonify({'error': 'Изображения слайдов не были созданы'}), 500

        slides = []
        for i, img_path in enumerate(image_paths):
            if i >= len(notes):
                break
            with open(img_path, 'rb') as img_file:
                img_base64 = base64.b64encode(img_file.read()).decode('ascii')

            slides.append({
                'slide_index': i + 1,
                'notes': notes[i],
                'image_base64': img_base64
            })

        response = jsonify({
            'filename': file.filename,
            'total_slides': len(slides),
            'slides': slides
        })
        response.headers.add('Access-Control-Allow-Origin', '*')
        response.headers.add('Access-Control-Allow-Headers', 'Content-Type')
        return response


if __name__ == '__main__':
    app.run(host='0.0.0.0', port=8081, debug=True)
